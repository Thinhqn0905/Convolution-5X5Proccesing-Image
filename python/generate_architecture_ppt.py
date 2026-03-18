#!/usr/bin/env python3
"""Generate architecture slide deck for convolution_fpga seminar.

Usage:
    python python/generate_architecture_ppt.py --out docs/architecture_convolution_fpga.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for idx, line in enumerate(bullets):
        if idx == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0


def _add_flow_slide(prs):
    from pptx.util import Inches  # type: ignore[import-not-found]

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Dataflow Architecture"

    labels = [
        "AXI-Stream In",
        "line_buffer_4",
        "kernel_loader",
        "mac_array_25x3",
        "AXI-Stream Out",
    ]

    x = Inches(0.4)
    y = Inches(2.0)
    w = Inches(2.2)
    h = Inches(0.9)
    gap = Inches(0.2)

    for i, label in enumerate(labels):
        rect = slide.shapes.add_shape(1, x + i * (w + gap), y, w, h)  # MSO_SHAPE.RECTANGLE=1
        rect.text = label

    ctrl = slide.shapes.add_shape(1, Inches(3.0), Inches(3.4), Inches(4.1), Inches(0.9))
    ctrl.text = "AXI-Lite Control: axi_lite_kernel_ctrl"


def build_presentation(out_path: Path) -> None:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required. Install with: pip install -r python/requirements.txt"
        ) from exc

    prs = Presentation()

    _add_title_slide(
        prs,
        title="RGB 5x5 Convolution Engine on FPGA",
        subtitle="Architecture, Verification, and Timing Closure Summary",
    )

    _add_bullets_slide(
        prs,
        "System Goal",
        [
            "Runtime-programmable 5x5 RGB convolution on streaming input",
            "AXI-Stream datapath + AXI-Lite kernel control",
            "Simulation-first signoff before board deployment",
        ],
    )

    _add_flow_slide(prs)

    _add_bullets_slide(
        prs,
        "Numeric and Pipeline Details",
        [
            "Pixel: RGB888, Coeff: signed 16-bit, Accumulator: signed 48-bit",
            "Fixed-point normalization: KERNEL_Q = 8",
            "MAC pipeline staged for timing closure and stable throughput",
            "Output valid delayed by internal MAC pipeline stages",
        ],
    )

    _add_bullets_slide(
        prs,
        "Verification and Timing Status",
        [
            "Regression PASS: identity5, gaussian5, sharpen5, emboss5, laplacian5",
            "Timing clean at 40 MHz (25 ns)",
            "Stretch points 50/60 MHz still pending closure",
            "Next gate: AXI backpressure hardening + SAIF confidence uplift",
        ],
    )

    _add_bullets_slide(
        prs,
        "Roadmap to Final Demo",
        [
            "Gate F hardening: randomized backpressure and interface stress",
            "Power confidence uplift with realistic SAIF activity",
            "Reproducible benchmark matrix (resolution x kernel)",
            "Seminar package: waveform proofs, timing table, live demo",
        ],
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate architecture pptx.")
    parser.add_argument(
        "--out",
        type=Path,
        default=Path("docs/architecture_convolution_fpga.pptx"),
        help="Output pptx path",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    build_presentation(args.out)
    print(f"Generated: {args.out}")


if __name__ == "__main__":
    main()
